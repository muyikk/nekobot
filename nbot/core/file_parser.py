"""
本地文件解析模块
支持 PDF、DOCX、PPT、Excel、TXT 等多种文件格式的文本提取
"""
import logging
import os
from typing import Dict, Any, Optional

_log = logging.getLogger(__name__)

# 编程语言文件扩展名（使用纯文本解析）
CODE_EXTENSIONS = {
    '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.c', '.cpp', '.h', '.hpp',
    '.cs', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.scala', '.r',
    '.lua', '.sh', '.bash', '.zsh', '.ps1', '.bat', '.cmd', '.pl', '.perl',
    '.sql', '.html', '.htm', '.css', '.scss', '.sass', '.less', '.xml', '.yaml',
    '.yml', '.toml', '.ini', '.cfg', '.conf', '.json', '.md', '.rst', '.tex',
    '.vue', '.svelte', '.jsx', '.tsx', '.dart', '.groovy', '.gradle', '.makefile',
    '.dockerfile', '.gitignore', '.env', '.properties'
}

# 支持的文档格式及对应解析器
DOCUMENT_EXTENSIONS = {
    '.pdf': 'pdf',
    '.docx': 'docx',
    '.doc': 'docx',
    '.pptx': 'pptx',
    '.ppt': 'pptx',
    '.xlsx': 'excel',
    '.xls': 'excel',
    '.txt': 'txt',
    '.csv': 'csv',
    '.md': 'markdown',
    '.rtf': 'rtf'
}


class FileParser:
    """文件解析器类"""
    
    @staticmethod
    def is_code_file(filename: str) -> bool:
        """判断是否为代码文件"""
        _, ext = os.path.splitext(filename.lower())
        return ext in CODE_EXTENSIONS
    
    @staticmethod
    def get_parser_type(filename: str) -> str:
        """获取文件的解析类型
        
        Returns:
            'code' - 代码文件，使用文本解析
            'pdf' - PDF文件
            'docx' - Word文档
            'pptx' - PowerPoint文档
            'excel' - Excel文档
            'txt' - 纯文本
            'unsupported' - 不支持的文件类型
        """
        _, ext = os.path.splitext(filename.lower())
        
        # 代码文件使用文本解析
        if ext in CODE_EXTENSIONS:
            return 'code'
        
        # 文档格式
        if ext in DOCUMENT_EXTENSIONS:
            return DOCUMENT_EXTENSIONS[ext]
        
        return 'unsupported'
    
    @staticmethod
    def parse_file(file_path: str, filename: str = None, max_chars: int = 50000) -> Optional[Dict[str, Any]]:
        """解析文件并提取文本内容
        
        Args:
            file_path: 文件完整路径
            filename: 文件名（可选）
            max_chars: 最大提取字符数
            
        Returns:
            包含解析结果的字典，失败返回 None
        """
        if not os.path.exists(file_path):
            _log.error(f"文件不存在: {file_path}")
            return None
        
        if filename is None:
            filename = os.path.basename(file_path)
        
        parser_type = FileParser.get_parser_type(filename)
        
        if parser_type == 'unsupported':
            _log.warning(f"不支持的文件类型: {filename}")
            return {
                'success': False,
                'error': f'不支持的文件类型: {os.path.splitext(filename)[1]}',
                'filename': filename,
                'type': 'unsupported'
            }
        
        try:
            if parser_type in ['code', 'txt', 'markdown', 'rtf', 'csv']:
                return FileParser._parse_text_file(file_path, filename, max_chars)
            elif parser_type == 'pdf':
                return FileParser._parse_pdf(file_path, filename, max_chars)
            elif parser_type == 'docx':
                return FileParser._parse_docx(file_path, filename, max_chars)
            elif parser_type == 'pptx':
                return FileParser._parse_pptx(file_path, filename, max_chars)
            elif parser_type == 'excel':
                return FileParser._parse_excel(file_path, filename, max_chars)
            else:
                return {
                    'success': False,
                    'error': f'未知的解析类型: {parser_type}',
                    'filename': filename
                }
        except Exception as e:
            _log.error(f"解析文件失败 {filename}: {e}")
            return {
                'success': False,
                'error': str(e),
                'filename': filename
            }
    
    @staticmethod
    def _parse_text_file(file_path: str, filename: str, max_chars: int) -> Dict[str, Any]:
        """解析文本文件（代码、txt、markdown等）"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        original_length = len(content)
        if len(content) > max_chars:
            content = content[:max_chars]
            truncated = True
        else:
            truncated = False
        
        return {
            'success': True,
            'content': content,
            'filename': filename,
            'type': 'text',
            'original_length': original_length,
            'extracted_length': len(content),
            'truncated': truncated,
            'file_size': os.path.getsize(file_path)
        }
    
    @staticmethod
    def _parse_pdf(file_path: str, filename: str, max_chars: int) -> Dict[str, Any]:
        """解析 PDF 文件"""
        try:
            import pdfplumber
            
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"[第 {i+1}/{total_pages} 页]\n{page_text}")
                    
                    # 提前截断以提高性能
                    if sum(len(t) for t in text_parts) > max_chars:
                        break
            
            content = '\n\n'.join(text_parts)
            original_length = len(content)
            if len(content) > max_chars:
                content = content[:max_chars]
                truncated = True
            else:
                truncated = False
            
            return {
                'success': True,
                'content': content,
                'filename': filename,
                'type': 'pdf',
                'pages': total_pages,
                'original_length': original_length,
                'extracted_length': len(content),
                'truncated': truncated,
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            _log.warning("pdfplumber 未安装，尝试使用 PyPDF2")
            return FileParser._parse_pdf_pypdf2(file_path, filename, max_chars)
    
    @staticmethod
    def _parse_pdf_pypdf2(file_path: str, filename: str, max_chars: int) -> Dict[str, Any]:
        """使用 PyPDF2 解析 PDF（备用方案）"""
        try:
            from PyPDF2 import PdfReader
            
            text_parts = []
            reader = PdfReader(file_path)
            total_pages = len(reader.pages)
            
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[第 {i+1}/{total_pages} 页]\n{page_text}")
                
                if sum(len(t) for t in text_parts) > max_chars:
                    break
            
            content = '\n\n'.join(text_parts)
            original_length = len(content)
            if len(content) > max_chars:
                content = content[:max_chars]
                truncated = True
            else:
                truncated = False
            
            return {
                'success': True,
                'content': content,
                'filename': filename,
                'type': 'pdf',
                'pages': total_pages,
                'original_length': original_length,
                'extracted_length': len(content),
                'truncated': truncated,
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            return {
                'success': False,
                'error': 'PDF 解析失败：未安装 pdfplumber 或 PyPDF2',
                'filename': filename
            }
    
    @staticmethod
    def _parse_docx(file_path: str, filename: str, max_chars: int) -> Dict[str, Any]:
        """解析 Word 文档"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs = []
            
            # 提取段落
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # 提取表格
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        tables_text.append(row_text)
                tables_text.append('')  # 表格之间添加空行
            
            content = '\n'.join(paragraphs)
            if tables_text:
                content += '\n\n【表格内容】\n' + '\n'.join(tables_text)
            
            original_length = len(content)
            if len(content) > max_chars:
                content = content[:max_chars]
                truncated = True
            else:
                truncated = False
            
            return {
                'success': True,
                'content': content,
                'filename': filename,
                'type': 'docx',
                'paragraphs': len(paragraphs),
                'tables': len(doc.tables),
                'original_length': original_length,
                'extracted_length': len(content),
                'truncated': truncated,
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            return {
                'success': False,
                'error': 'Word 文档解析失败：未安装 python-docx',
                'filename': filename
            }
    
    @staticmethod
    def _parse_pptx(file_path: str, filename: str, max_chars: int) -> Dict[str, Any]:
        """解析 PowerPoint 文档"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"[幻灯片 {slide_num}/{len(prs.slides)}]"]
                
                # 提取标题
                if slide.shapes.title:
                    slide_parts.append(f"标题: {slide.shapes.title.text}")
                
                # 提取文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text and text != slide.shapes.title.text if slide.shapes.title else True:
                            slide_parts.append(text)
                
                slides_text.append('\n'.join(slide_parts))
                
                # 提前截断
                if sum(len(t) for t in slides_text) > max_chars:
                    break
            
            content = '\n\n'.join(slides_text)
            original_length = len(content)
            if len(content) > max_chars:
                content = content[:max_chars]
                truncated = True
            else:
                truncated = False
            
            return {
                'success': True,
                'content': content,
                'filename': filename,
                'type': 'pptx',
                'slides': len(prs.slides),
                'original_length': original_length,
                'extracted_length': len(content),
                'truncated': truncated,
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            return {
                'success': False,
                'error': 'PowerPoint 解析失败：未安装 python-pptx',
                'filename': filename
            }
    
    @staticmethod
    def _parse_excel(file_path: str, filename: str, max_chars: int) -> Dict[str, Any]:
        """解析 Excel 文档"""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheets_text = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_parts = [f"[工作表: {sheet_name}]"]
                
                rows_data = []
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    # 跳过空行
                    if any(cell is not None for cell in row):
                        row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                        rows_data.append(row_text)
                        row_count += 1
                        
                        # 限制每张表的最大行数
                        if row_count > 1000:
                            rows_data.append('...(超过1000行，已截断)')
                            break
                
                if rows_data:
                    sheet_parts.extend(rows_data)
                
                sheets_text.append('\n'.join(sheet_parts))
                
                # 提前截断
                if sum(len(t) for t in sheets_text) > max_chars:
                    break
            
            wb.close()
            
            content = '\n\n'.join(sheets_text)
            original_length = len(content)
            if len(content) > max_chars:
                content = content[:max_chars]
                truncated = True
            else:
                truncated = False
            
            return {
                'success': True,
                'content': content,
                'filename': filename,
                'type': 'excel',
                'sheets': len(wb.sheetnames),
                'sheet_names': wb.sheetnames if hasattr(wb, 'sheetnames') else [],
                'original_length': original_length,
                'extracted_length': len(content),
                'truncated': truncated,
                'file_size': os.path.getsize(file_path)
            }
        except ImportError:
            return {
                'success': False,
                'error': 'Excel 解析失败：未安装 openpyxl',
                'filename': filename
            }
    
    @staticmethod
    def get_file_metadata(file_path: str, filename: str = None) -> Dict[str, Any]:
        """获取文件元数据（不解析内容）
        
        Args:
            file_path: 文件完整路径
            filename: 文件名（可选）
            
        Returns:
            文件元数据字典
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'error': '文件不存在'
            }
        
        if filename is None:
            filename = os.path.basename(file_path)
        
        _, ext = os.path.splitext(filename.lower())
        file_size = os.path.getsize(file_path)
        
        # 格式化文件大小
        if file_size < 1024:
            size_str = f"{file_size} B"
        elif file_size < 1024 * 1024:
            size_str = f"{file_size / 1024:.1f} KB"
        else:
            size_str = f"{file_size / (1024 * 1024):.1f} MB"
        
        parser_type = FileParser.get_parser_type(filename)
        
        # 获取页数/工作表数等信息（不解析全文）
        extra_info = {}
        
        try:
            if parser_type == 'pdf':
                try:
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        extra_info['pages'] = len(pdf.pages)
                except:
                    try:
                        from PyPDF2 import PdfReader
                        reader = PdfReader(file_path)
                        extra_info['pages'] = len(reader.pages)
                    except:
                        pass
            
            elif parser_type == 'docx':
                try:
                    from docx import Document
                    doc = Document(file_path)
                    extra_info['paragraphs'] = len([p for p in doc.paragraphs if p.text.strip()])
                    extra_info['tables'] = len(doc.tables)
                except:
                    pass
            
            elif parser_type == 'pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    extra_info['slides'] = len(prs.slides)
                except:
                    pass
            
            elif parser_type == 'excel':
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path, read_only=True)
                    extra_info['sheets'] = len(wb.sheetnames)
                    extra_info['sheet_names'] = wb.sheetnames[:5]  # 只显示前5个工作表名
                    if len(wb.sheetnames) > 5:
                        extra_info['sheet_names'].append(f'... 共 {len(wb.sheetnames)} 个')
                    wb.close()
                except:
                    pass
        except Exception as e:
            _log.warning(f"获取文件元数据时出错: {e}")
        
        return {
            'success': True,
            'filename': filename,
            'type': parser_type,
            'extension': ext,
            'size': file_size,
            'size_str': size_str,
            'can_parse': parser_type != 'unsupported',
            **extra_info
        }


# 全局解析器实例
file_parser = FileParser()
